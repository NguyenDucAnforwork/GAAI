"""GAIA attachment resolution + text extraction.

Resolution order:
  1. Official gated dataset `gaia-benchmark/GAIA` via HF token (3 of the
     supplied tokens have access).
  2. Public community mirror `datasets/asteriadyt/2023`.
  3. Scoring API /files/<task_id> (usually disabled in this env, last resort).
"""
import os
import io

import requests
from dotenv import load_dotenv, find_dotenv

load_dotenv(find_dotenv(), override=True)

# Tokens with confirmed GAIA access first (cbg/hust/forwork); geminipro lacks it.
_HF_TOKENS = [t for t in (
    os.getenv("HF_TOKEN_cbg"),
    os.getenv("HF_TOKEN_hust"),
    os.getenv("HF_TOKEN_forwork"),
    os.getenv("HF_TOKEN_geminipro"),
) if t]

_MIRROR = "https://huggingface.co/datasets/asteriadyt/2023/resolve/main/validation"
_SCORING = "https://agents-course-unit4-scoring.hf.space"

IMAGE_EXT = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")
AUDIO_EXT = (".mp3", ".wav", ".m4a", ".flac", ".ogg")
TEXTUAL_EXT = (".txt", ".csv", ".json", ".jsonl", ".xml", ".md", ".py", ".tsv")


def fetch_file(file_name: str, task_id: str = None) -> bytes:
    """Return the raw bytes of a GAIA attachment, trying official → mirror → API."""
    # 1) Official gated dataset
    try:
        from huggingface_hub import hf_hub_download
        for tok in _HF_TOKENS:
            try:
                p = hf_hub_download(
                    repo_id="gaia-benchmark/GAIA",
                    filename=f"2023/validation/{file_name}",
                    repo_type="dataset", token=tok,
                )
                with open(p, "rb") as f:
                    return f.read()
            except Exception:
                continue
    except Exception:
        pass
    # 2) Public mirror
    try:
        r = requests.get(f"{_MIRROR}/{file_name}",
                         headers={"User-Agent": "Mozilla/5.0"}, timeout=40)
        if r.status_code == 200 and len(r.content) > 50:
            return r.content
    except Exception:
        pass
    # 3) Scoring API (by task_id)
    if task_id:
        try:
            r = requests.get(f"{_SCORING}/files/{task_id}", timeout=30)
            if r.status_code == 200 and len(r.content) > 50:
                return r.content
        except Exception:
            pass
    raise RuntimeError(f"Could not fetch attachment: {file_name}")


def extract_text(file_name: str, content: bytes) -> str:
    """Extract plain text from a document attachment (empty if not text-like)."""
    fn = file_name.lower()
    if fn.endswith(".docx"):
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for t in doc.tables:
            for row in t.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    if fn.endswith((".xlsx", ".xls")):
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"[Sheet: {ws.title}]")
            # Values grid
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    out.append(" | ".join(cells))
            # Cell fill colours (many GAIA spreadsheets encode a map/grid in colour).
            colors = []
            for row in ws.iter_rows():
                for cell in row:
                    try:
                        fg = cell.fill.fgColor
                        rgb = getattr(fg, "rgb", None)
                        if rgb and isinstance(rgb, str) and rgb not in ("00000000",):
                            hexv = rgb[-6:].upper()
                            colors.append(f"{cell.coordinate}={hexv}")
                    except Exception:
                        continue
            if colors:
                out.append("[Cell fill colours (coordinate=RRGGBB)]")
                out.append(" ".join(colors))
        return "\n".join(out)
    if fn.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts.append(sh.text_frame.text.strip())
                if sh.has_table:
                    for r in sh.table.rows:
                        texts.append(" | ".join(c.text for c in r.cells))
            out.append(f"[Slide {i}] " + " || ".join(texts))
        return "\n".join(out)
    if fn.endswith(".pdf"):
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        return "\n".join((p.extract_text() or "") for p in reader.pages)
    if fn.endswith(TEXTUAL_EXT):
        return content.decode("utf-8", "ignore")
    return ""


def mime_for(file_name: str) -> str:
    fn = file_name.lower()
    if fn.endswith(".png"):
        return "image/png"
    if fn.endswith((".jpg", ".jpeg")):
        return "image/jpeg"
    if fn.endswith(".gif"):
        return "image/gif"
    if fn.endswith(".webp"):
        return "image/webp"
    if fn.endswith(".mp3"):
        return "audio/mp3"
    if fn.endswith(".wav"):
        return "audio/wav"
    if fn.endswith((".m4a", ".aac")):
        return "audio/aac"
    if fn.endswith(".flac"):
        return "audio/flac"
    if fn.endswith(".ogg"):
        return "audio/ogg"
    return "application/octet-stream"
